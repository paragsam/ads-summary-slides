"""
Generate a PPTX slide deck from TikTok Ads JSON data.

Produces:
  Slide 1  — Performance By Objective (one column per campaign objective)
  Slide 2+ — Performance By Creative, one slide per objective, top 5 by spend

Usage:
  python generate_pptx.py \
    --campaigns output/report_campaigns.json \
    --creatives output/report_creatives.json \
    --template templates/Performance_By_Objective_template.pptx \
    --start 2026-05-10 \
    --end 2026-05-16 \
    --out output/report.pptx
"""

from __future__ import annotations

import argparse
import copy
import json
import sys
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


# ---------------------------------------------------------------------------
# Formatting helpers
# ---------------------------------------------------------------------------

def fmt_currency(value: str) -> str:
    try:
        return f"${float(value):,.2f}"
    except (ValueError, TypeError):
        return value


def fmt_number(value: str) -> str:
    try:
        return f"{int(float(value)):,}"
    except (ValueError, TypeError):
        return value


def fmt_pct(value: str) -> str:
    """Add % if not already present."""
    v = str(value).strip()
    if v.endswith("%"):
        return v
    try:
        return f"{float(v):.2f}%"
    except (ValueError, TypeError):
        return v


def fmt_date_range(start: str, end: str) -> str:
    """Format '2026-05-10' / '2026-05-16' → '5/10-5/16'."""
    s = datetime.strptime(start, "%Y-%m-%d")
    e = datetime.strptime(end, "%Y-%m-%d")
    return f"{s.month}/{s.day}-{e.month}/{e.day}"


# ---------------------------------------------------------------------------
# Slide helpers
# ---------------------------------------------------------------------------

def replace_text_in_slide(slide, old: str, new: str) -> None:
    """Replace exact text in all text frames on a slide."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if old in run.text:
                    run.text = run.text.replace(old, new)


def set_cell_text(cell, text: str) -> None:
    """Write text into a table cell, preserving font formatting of the first run."""
    tf = cell.text_frame
    if not tf.paragraphs:
        return
    para = tf.paragraphs[0]
    if para.runs:
        first_run = para.runs[0]
        first_run.text = text
        for extra in para.runs[1:]:
            extra.text = ""
    else:
        run = para.add_run()
        run.text = text


def fill_table_rows(table, data_rows: list[list[str]], start_row: int = 1) -> None:
    """Write data_rows into the table starting at start_row (0-indexed)."""
    for r_idx, row_values in enumerate(data_rows):
        tbl_row_idx = start_row + r_idx
        if tbl_row_idx >= len(table.rows):
            break
        for c_idx, value in enumerate(row_values):
            if c_idx >= len(table.columns):
                break
            set_cell_text(table.rows[tbl_row_idx].cells[c_idx], value)


def ensure_table_rows(table, needed: int, header_rows: int = 1) -> None:
    """Clone the last data row until the table has enough rows for `needed` data rows."""
    current_data_rows = len(table.rows) - header_rows
    if current_data_rows >= needed:
        return
    last_tr = table.rows[len(table.rows) - 1]._tr
    tbl_element = last_tr.getparent()
    for _ in range(needed - current_data_rows):
        new_tr = copy.deepcopy(last_tr)
        # Clear text in the cloned row
        for tc in new_tr.findall(f".//{{{NS}}}tc"):
            for r in tc.findall(f".//{{{NS}}}r"):
                t = r.find(f"{{{NS}}}t")
                if t is not None:
                    t.text = ""
        tbl_element.append(new_tr)


def clone_slide(prs: Presentation, source_index: int) -> object:
    """Duplicate a slide by index and append the copy to the presentation."""
    source = prs.slides[source_index]
    layout = source.slide_layout
    new_slide = prs.slides.add_slide(layout)

    # Replace the new slide's shape tree content with a deep copy of the source
    sp_tree = new_slide.shapes._spTree
    for child in list(sp_tree):
        sp_tree.remove(child)
    for child in source.shapes._spTree:
        sp_tree.append(copy.deepcopy(child))

    return new_slide


def find_table(slide) -> object | None:
    for shape in slide.shapes:
        if shape.has_table:
            return shape.table
    return None


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

OBJECTIVE_LABEL = {
    "REACH": "Reach",
    "TRAFFIC": "Traffic",
    "VIDEO_VIEWS": "Video Views",
    "CONVERSIONS": "Conversions",
    "APP_PROMOTION": "App Promotion",
    "LEAD_GENERATION": "Lead Generation",
    "CATALOG_SALES": "Catalog Sales",
    "COMMUNITY_INTERACTION": "Community",
}


def objective_label(raw: str) -> str:
    return OBJECTIVE_LABEL.get(raw, raw.replace("_", " ").title())


def build_objective_slide(slide, campaigns: list[dict], date_label: str) -> None:
    """Fill Slide 1 — Performance By Objective."""
    table = find_table(slide)
    if table is None:
        return

    replace_text_in_slide(slide, "11/17-11/30", date_label)

    # Column headers (row 0, cols 1+)
    for col_idx, campaign in enumerate(campaigns, start=1):
        if col_idx >= len(table.columns):
            break
        set_cell_text(table.rows[0].cells[col_idx], objective_label(campaign["objective_type"]))

    # Clear remaining header columns if fewer campaigns than columns
    for col_idx in range(len(campaigns) + 1, len(table.columns)):
        set_cell_text(table.rows[0].cells[col_idx], "")

    # Data rows
    metric_rows = [
        ("Cost",             lambda c: fmt_currency(c["spend"])),
        ("Impressions",      lambda c: fmt_number(c["impressions"])),
        ("CPM",              lambda c: fmt_currency(c["cpm"])),
        ("Reach",            lambda c: fmt_number(c["reach"])),
        ("Clicks",           lambda c: fmt_number(c["clicks"])),
        ("CTR",              lambda c: fmt_pct(c["ctr"])),
        ("CPC",              lambda c: fmt_currency(c["cpc"])),
        ("VTR",              lambda c: fmt_pct(c["vtr"])),
        ("6-Sec VTR",        lambda c: fmt_pct(c["vtr_6s"])),
        ("Engagement Rate",  lambda c: fmt_pct(c["engagement_rate"])),
    ]

    for row_idx, (_, extractor) in enumerate(metric_rows, start=1):
        if row_idx >= len(table.rows):
            break
        for col_idx, campaign in enumerate(campaigns, start=1):
            if col_idx >= len(table.columns):
                break
            set_cell_text(table.rows[row_idx].cells[col_idx], extractor(campaign))
        # Clear remaining columns
        for col_idx in range(len(campaigns) + 1, len(table.columns)):
            set_cell_text(table.rows[row_idx].cells[col_idx], "")


def build_creative_slide(slide, objective: str, creatives: list[dict], date_label: str) -> None:
    """Fill a creative slide — Performance By Creative for one objective."""
    replace_text_in_slide(slide, "Performance By Creative",
                          f"Performance By Creative — {objective_label(objective)}")
    replace_text_in_slide(slide, "11/17-11/30", date_label)

    table = find_table(slide)
    if table is None:
        return

    top = sorted(creatives, key=lambda c: float(c["spend"]), reverse=True)[:5]

    # Ensure the table has enough rows for up to 5 creatives
    ensure_table_rows(table, needed=len(top), header_rows=1)

    # Clear all existing data rows first
    for r in range(1, len(table.rows)):
        for c in range(len(table.columns)):
            set_cell_text(table.rows[r].cells[c], "")

    # Fill in the top creatives
    for r_idx, creative in enumerate(top, start=1):
        row_data = [
            creative["ad_name"],
            fmt_currency(creative["spend"]),
            fmt_number(creative["impressions"]),
            fmt_currency(creative["cpm"]),
            fmt_pct(creative["ctr"]),
            fmt_currency(creative["cpc"]),
            fmt_pct(creative["vtr_6s"]),
            fmt_pct(creative["vtr"]),
            fmt_pct(creative["engagement_rate"]),
        ]
        for c_idx, value in enumerate(row_data):
            if c_idx < len(table.columns):
                set_cell_text(table.rows[r_idx].cells[c_idx], value)


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

TEMPLATE_DIR = Path(__file__).parent / "templates"
MAX_CAMPAIGNS = 5


def resolve_template(campaign_count: int) -> Path:
    """Pick the right template based on number of campaigns."""
    count = max(1, min(campaign_count, MAX_CAMPAIGNS))
    path = TEMPLATE_DIR / f"Brand_Objective_{count}_campaign.pptx"
    if not path.exists():
        print(
            json.dumps({
                "error": f"Template not found for {count} campaign(s): {path}",
                "detail": (
                    f"Expected a file named 'Brand_Objective_{count}_campaign.pptx' "
                    f"in the templates/ directory."
                ),
            }),
            file=sys.stderr,
        )
        sys.exit(1)
    return path


def parse_args(argv: list[str] | None = None) -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Generate a PPTX deck from TikTok Ads JSON data."
    )
    parser.add_argument("--campaigns", required=True, help="Path to _campaigns.json")
    parser.add_argument("--creatives", required=True, help="Path to _creatives.json")
    parser.add_argument("--start", required=True, help="Start date YYYY-MM-DD")
    parser.add_argument("--end", required=True, help="End date YYYY-MM-DD")
    parser.add_argument("--out", required=True, help="Output .pptx path")
    return parser.parse_args(argv)


def main(argv: list[str] | None = None) -> None:
    args = parse_args(argv)

    campaigns_path = Path(args.campaigns)
    creatives_path = Path(args.creatives)
    out_path = Path(args.out)

    for p in (campaigns_path, creatives_path):
        if not p.exists():
            print(json.dumps({"error": f"File not found: {p}"}), file=sys.stderr)
            sys.exit(1)

    campaigns = json.loads(campaigns_path.read_text())
    creatives = json.loads(creatives_path.read_text())
    date_label = fmt_date_range(args.start, args.end)

    template_path = resolve_template(len(campaigns))
    print(f"Using template: {template_path.name} ({len(campaigns)} campaign(s))")
    prs = Presentation(str(template_path))

    # --- Slide 1: Performance By Objective ---
    print("Building Slide 1: Performance By Objective...")
    build_objective_slide(prs.slides[0], campaigns, date_label)

    # --- Slides 2+: Performance By Creative, one per objective ---
    objectives = list(dict.fromkeys(c["objective_type"] for c in campaigns))
    for i, objective in enumerate(objectives):
        obj_creatives = [c for c in creatives if c["objective_type"] == objective]
        print(f"Building creative slide for objective: {objective_label(objective)} "
              f"({len(obj_creatives)} creatives)...")

        if i == 0:
            # Use the existing Slide 2 as-is for the first objective
            creative_slide = prs.slides[1]
        else:
            # Clone Slide 2 for subsequent objectives
            creative_slide = clone_slide(prs, 1)

        build_creative_slide(creative_slide, objective, obj_creatives, date_label)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"Saved: {out_path.resolve()}")


if __name__ == "__main__":
    main()
